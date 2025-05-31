import os
import shutil
import mimetypes
import fitz  # PyMuPDF
import docx
import pandas as pd
from pptx import Presentation
from PIL import Image
import pytesseract
import json

from google import genai
from google.genai import types

# ==== Setup Gemini Client ====
client = genai.Client(api_key='INSERT_YOUR_API_KEY')
model = "gemini-2.5-flash-preview-05-20"

generate_content_config = types.GenerateContentConfig(
    response_mime_type="application/json",
    response_schema=genai.types.Schema(
        type=genai.types.Type.OBJECT,
        required=["content", "label", "filename"],
        properties={
            "content": genai.types.Schema(type=genai.types.Type.STRING),
            "label": genai.types.Schema(type=genai.types.Type.STRING),
            "filename": genai.types.Schema(type=genai.types.Type.STRING),
        },
    ),
)

def get_label_from_gemini(content, filename):
    contents = [
        types.Content(
            role="user",
            parts=[types.Part(text=f"{content}\n\nfilename : {filename}")]
        )
    ]
    response = client.models.generate_content(
        model=model,
        contents=contents,
        config=generate_content_config
    )
    return response.text

def read_file_content(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".pdf":
            return read_pdf(file_path)
        elif ext == ".docx":
            return read_docx(file_path)
        elif ext == ".csv":
            return read_csv(file_path)
        elif ext in [".xls", ".xlsx"]:
            return read_excel(file_path)
        elif ext == ".pptx":
            return read_pptx(file_path)
        elif ext in [".jpg", ".jpeg", ".png", ".bmp", ".tiff"]:
            return read_image_text(file_path)
        elif ext in [".txt", ".py", ".json", ".html", ".xml", ".md"]:
            return read_text_file(file_path)
        else:
            return try_generic_read(file_path)
    except Exception as e:
        print(f"❌ Error reading {file_path}: {e}")
        return None

# === Format-specific readers ===

def read_text_file(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def read_pdf(file_path):
    text = ""
    doc = fitz.open(file_path)
    for page in doc:
        text += page.get_text()
    return text

def read_docx(file_path):
    doc = docx.Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs)

def read_csv(file_path):
    df = pd.read_csv(file_path)
    return df.to_string(index=False)

def read_excel(file_path):
    df = pd.read_excel(file_path)
    return df.to_string(index=False)

def read_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def read_image_text(file_path):
    img = Image.open(file_path)
    img = img.convert("RGB")  # Ensure compatibility with Tesseract
    return pytesseract.image_to_string(img)

def try_generic_read(file_path):
    try:
        with open(file_path, "rb") as f:
            return f.read().decode("utf-8", errors="ignore")
    except:
        return "[Unreadable binary file]"

def organize_files_by_label(root_dir):
    for folderpath, _, filenames in os.walk(root_dir):
        for fname in filenames:
            fpath = os.path.join(folderpath, fname)
            print(f"\n📂 Processing: {fpath}")
            content = read_file_content(fpath)
            if not content:
                continue

            try:
                json_output = get_label_from_gemini(content, fname)
                print(f"✅ Gemini Output: {json_output}")
                data = json.loads(json_output)
                label = data["label"]
            except Exception as e:
                print(f"❌ Error calling Gemini: {e}")
                continue

            parent_folder = folderpath
            target_folder = os.path.join(parent_folder, label)

            # Skip if already in correct subfolder
            if os.path.normpath(os.path.dirname(fpath)) == os.path.normpath(target_folder):
                print("📎 File already in correct folder. Skipping.")
                continue

            # Ask for approval before moving
            print(f"📝 Suggested move: {fname} → {label}/")
            choice = input("Do you want to move this file? [y/N]: ").strip().lower()
            if choice != "y":
                print("❌ Skipped.")
                continue

            os.makedirs(target_folder, exist_ok=True)

            try:
                shutil.move(fpath, os.path.join(target_folder, fname))
                print(f"📁 Moved to {target_folder}")
            except Exception as e:
                print(f"❌ Failed to move file: {e}")

# ===== Run the Script =====
if __name__ == "__main__":
    base_path = input("Enter the root folder path: ").strip()
    if os.path.exists(base_path):
        organize_files_by_label(base_path)
    else:
        print("❌ Path not found.")
